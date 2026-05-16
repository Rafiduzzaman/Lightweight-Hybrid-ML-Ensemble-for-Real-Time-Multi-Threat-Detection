import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "results" / "presentation_ensemble_20260515.pptx"


def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_text_slide(prs, heading, lines):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
            p.text = line
        else:
            p = body.add_paragraph()
            p.text = line
        p.level = 0


def add_image_slide(prs, heading, img_path, width_inches=8):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = heading
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.4), width=Inches(width_inches))


def main():
    prs = Presentation()

    # Title
    add_title_slide(prs, "Lightweight Hybrid ML Ensemble for Real-Time Multi-Threat Detection",
                    "Summary of methods, results, and showcase")

    # Extract abstract and key lines from the paper
    tex_path = ROOT / "ieee_thesis_paper.tex"
    abstract = ""
    if tex_path.exists():
        text = tex_path.read_text(encoding="utf-8")
        # crude extraction
        if "\\begin{abstract}" in text and "\\end{abstract}" in text:
            abstract = text.split("\\begin{abstract}")[1].split("\\end{abstract}")[0].strip()

    if abstract:
        lines = [l.strip() for l in abstract.splitlines() if l.strip()][:6]
        add_text_slide(prs, "Abstract", lines)

    # Methods slide
    methods = [
        "Weighted soft-voting ensemble: KNN (0.2), RF (0.4), XGBoost (0.4)",
        "Preprocessing: numeric features, StandardScaler",
        "Feature selection: SelectKBest (k=30)",
        "Imbalance handling: SMOTE (k=5)",
        "Evaluation: Accuracy, Precision, Recall, F1, ROC-AUC; within- and cross-dataset"
    ]
    add_text_slide(prs, "Methodology", methods)

    # Datasets slide
    datasets = [
        "CIC-IDS-2017: ~2.8M records (80 features)",
        "UNSW-NB15: ~447K records (49 features)",
        "CIC-IDS-2018: sampled 90K records (80 features)"
    ]
    add_text_slide(prs, "Datasets", datasets)

    # Results: read summary CSV and add a concise table slide per dataset
    tables_dir = ROOT / "results_colab" / "tables"
    all_models_csv = tables_dir / "all_models_comparison_20260513_111922.csv"
    if all_models_csv.exists():
        df = pd.read_csv(all_models_csv)
        for ds in ["CIC-IDS-2017", "UNSW-NB15", "CIC-IDS-2018"]:
            sub = df[df["Dataset"] == ds].sort_values("Accuracy (%)", ascending=False)
            top = sub.head(4)
            lines = [f"{r.Model}: Acc={r['Accuracy (%)']:.2f}%, F1={r['F1-Score']:.3f}, ROC-AUC={r['ROC-AUC']:.3f}" for _, r in top.iterrows()]
            add_text_slide(prs, f"Top Models — {ds}", lines)

    # Ablation and cross-dataset
    ablation_csv = tables_dir / "ablation_study_20260513_111922.csv"
    if ablation_csv.exists():
        ab = pd.read_csv(ablation_csv)
        lines = []
        first_col = ab.columns[0]
        for _, r in ab.iterrows():
            lines.append(f"{r[first_col]}: Acc={r['Accuracy (%)']:.2f}%")
        add_text_slide(prs, "Ablation Study", lines[:6])

    cross_csv = tables_dir / "cross_dataset_evaluation_20260513_111922.csv"
    if cross_csv.exists():
        cross = pd.read_csv(cross_csv)
        lines = [f"Train {r['Train Dataset']} → Test {r['Test Dataset']}: Acc={r['Accuracy (%)']:.2f}% (Drop={r['Drop Percentage']})" for _, r in cross.iterrows()]
        add_text_slide(prs, "Cross-Dataset Evaluation", lines[:6])

    # Add figure slides if available
    figs = ROOT / "results_colab" / "figures"
    fig_map = [
        ("ROC Curves (CIC-2017)", figs / "roc_curves_cic_2017_20260513_111922.png"),
        ("XGBoost Loss Curves", figs / "xgboost_loss_curves_cic_2017_20260513_111922.png"),
        ("Ensemble Confusion Matrix", figs / "confusion_matrix_ensemble_cic_2017_20260513_111922.png"),
    ]
    for title, p in fig_map:
        add_image_slide(prs, title, p)

    # Conclusions
    conclusions = [
        "Ensemble achieves strong within-dataset performance (e.g., 99.56% on CIC-IDS-2017)",
        "Feature selection and ensemble diversity are main contributors",
        "Cross-dataset transfer shows major degradation — adaptation needed for deployment",
        "Future work: domain adaptation, continual learning, more minority sampling"
    ]
    add_text_slide(prs, "Conclusions", conclusions)

    # Save
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved presentation to: {OUT}")


if __name__ == "__main__":
    main()
